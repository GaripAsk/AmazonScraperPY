import os
from curl_cffi import requests
from bs4 import BeautifulSoup
import re
from pptx import Presentation
from pptx.util import Inches
from pptx.util import Pt
from PIL import Image
from selenium import webdriver
from selenium.webdriver.common.by import By
from selenium.webdriver.support.ui import WebDriverWait
from selenium.webdriver.support import expected_conditions as EC
from selenium.webdriver.firefox.options import Options
from PIL import Image
import time
from tkinter import Tk
from tkinter.filedialog import askopenfilename
from pptx.dml.color import RGBColor

def extract_numbers(text):
    return ''.join(re.findall(r'\d', text))

def the_machine(url):

    product_info = {}

    stripped_product_info = {}

    filtered_info = {}

    new_output = {}

    try:
        link = requests.get(url, headers=headers, impersonate="chrome110")
        sayfa_parsel = BeautifulSoup(link.content, "html.parser")

        price = sayfa_parsel.find("span", class_="a-offscreen").string

    except AttributeError:
        print("An error occurred while collecting, trying second method")
        try:
            options = Options()
            options.add_argument('-headless')

            print("Starting the browser")

            driver = webdriver.Firefox(options=options)
            driver.maximize_window()

            print("Browser started")

            driver.get(url)

            print("link uploaded")

            time.sleep(2)
            try:
                cookie_accept_button = WebDriverWait(driver, 20).until(EC.element_to_be_clickable((By.ID, "sp-cc-accept")))
                if cookie_accept_button:
                    cookie_accept_button.click()
                    print("cookies accepted")
            except Exception as e:
                print(f"Cookie button couldn't find, error: {e}")
            
            time.sleep(5)
            page_source = driver.page_source

            sayfa_parsel = BeautifulSoup(page_source, "html.parser")
            print("Html code have been taken")
        except Exception as e:
            print("Html code couldn't taken")
    #price collecter
    try:
        price = sayfa_parsel.find("span", class_="a-offscreen").string
        print("Price information collected")
    except:
        try:
            price = sayfa_parsel.find("div", class_="a-spacing-top-mini").find("span", class_="a-offscreen")
            print("Price information collected")
        except Exception as e:
            price = "-"
            print(f"An error occurred while collecting the price, error {e}")

    #title collecter
    try:
        product_title = sayfa_parsel.find("span", id="productTitle").string
        print("Title information collected")
    except Exception as e:
        print(f"An error occurred while collecting the title information, error: {e}")
        product_title = "-"

    #review collecter
    try:
        review = sayfa_parsel.find("span", id="acrCustomerReviewText").string
        print("Review information collected")
    except:
        try:
            review = sayfa_parsel.select_one('span[data-hook="total-review-count"].a-size-base.a-color-secondary').string
            print("Review information collected")
        except:
            try:
                review = sayfa_parsel.find("span", data_hook="total-review-count").string
                print("Review information collected")
            except Exception as e:
                print(f"An error occurred while collecting the review, error: {e}")
                review = "-"
    if review == None:
        print("review is still method")
        review = "-"
    
    #side info taker
    time.sleep(2)
    try:
        li_tags = sayfa_parsel.find("div", id="detailBullets_feature_div").find_all("li")

        for li in li_tags:
            key_span = li.find('span', {'class': 'a-text-bold'})
            value_span = li.find_all('span')[-1]  
                    
            if key_span and value_span:
                key = key_span.string if key_span.string else key_span.text
                value = value_span.string if value_span.string else value_span.text
                        
                key = re.sub(r'\s+', ' ', key).strip().replace(':', '').replace('\u200f', '').replace('\u200e', '')
                value = re.sub(r'\s+', ' ', value).strip().replace(':', '').replace('\u200f', '').replace('\u200e', '')
                    
                product_info[key] = value
        print("Additional info collected")
    except:
        try:
            tr_tags = sayfa_parsel.find("div", class_="a-row a-spacing-top-base").find_all("tr")
            if tr_tags is None:
                tr_tags = sayfa_parsel.find('div', {'id': 'prodDetails', 'class': 'a-section'}).find_all("tr")
                raise ValueError("tr_tags is empty, trying another method.")
                

            for tr in tr_tags:
                key_span = tr.find('th', {'class': 'a-color-secondary a-size-base prodDetSectionEntry'})
                value_span = tr.find('td', {'class': 'a-size-base prodDetAttrValue'}) 
                        
                if key_span and value_span:
                    key = key_span.string if key_span.string else key_span.text
                    value = value_span.string if value_span.string else value_span.text
                            
                    key = re.sub(r'\s+', ' ', key).strip().replace(':', '').replace('\u200f', '').replace('\u200e', '')
                    value = re.sub(r'\s+', ' ', value).strip().replace(':', '').replace('\u200f', '').replace('\u200e', '')
                        
                    product_info[key] = value
            print("Additional info collected")
        except:
            try:
                tr_tags = sayfa_parsel.find("div", class_="a-row a-spacing-top-base").find_all("tr")
                if tr_tags is None:
                    tr_tags = sayfa_parsel.find('div', {'id': 'prodDetails', 'class': 'a-section'}).find_all("tr")
                    raise ValueError("tr_tags is empty, trying another method.")

                for tr in tr_tags:
                    key_span = tr.find('th', class_="a-color-secondary a-size-base prodDetSectionEntry")
                    value_span = tr.find('td', class_= "a-size-base prodDetAttrValue") 
                            
                    if key_span and value_span:
                        key = key_span.string if key_span.string else key_span.text
                        value = value_span.string if value_span.string else value_span.text
                                
                        key = re.sub(r'\s+', ' ', key).strip().replace(':', '').replace('\u200f', '').replace('\u200e', '')
                        value = re.sub(r'\s+', ' ', value).strip().replace(':', '').replace('\u200f', '').replace('\u200e', '')
                            
                        product_info[key] = value    
                        print("Additional info collected")            
            except Exception as e:
                print(f"An error occurred while collecting the side information, error: {e}")
                product_info = "-"

    #seller name collecter
    try:
        seller_name = sayfa_parsel.find("span", class_="a-size-small tabular-buybox-text-message").string
        print( "Seller name collected")
    except:
        try:
            seller_name_checker = sayfa_parsel.find("a", class_="a-link-normal", id="bylineInfo").string
            match_for_en1 = re.search("Visit the (.+?) Store", seller_name_checker)
            match_for_en2 = re.search("Brand: (.+?)", seller_name_checker)
            match_for_de1 = re.search("Besuche den (.+?)-Store", seller_name_checker)
            match_for_de2 = re.search("Marke: (.+?)", seller_name_checker)
            match_for_tr1 = re.search("(.+?) Store’u ziyaret edin", seller_name_checker)
            match_for_tr2 = re.search("Marka: (.+?)", seller_name_checker)
            if match_for_en1:
                seller_name = match_for_en1.group(1)
            elif match_for_en2:
                seller_name = match_for_en2.group(1)
            elif match_for_de1:
                seller_name = match_for_de1.group(1)
            elif match_for_de2:
                seller_name = match_for_de2.group(1)
            elif match_for_tr1:
                seller_name = match_for_tr1.group(1)
            elif match_for_tr2:
                seller_name = match_for_tr2.group(1)
            print( "Seller name collected")
        except Exception as e:
                seller_name = "-"
                print(f"Seller name could not be fetched, error: {e}")
            

    #ss collecter
    try:
        options = Options()
        options.add_argument('-headless')

        print("Starting the browser")

        driver = webdriver.Firefox(options=options)
        driver.maximize_window()

        print("Browser started")

        driver.get(url)

        print("link uploaded")

        time.sleep(2)
        try:
            cookie_accept_button = WebDriverWait(driver, 20).until(EC.element_to_be_clickable((By.ID, "sp-cc-accept")))
            if cookie_accept_button:
                cookie_accept_button.click()
                print("cookies accepted")
        except Exception as e:
            print(f"Cookie button couldn't find, error: {e}")

        print("Scrolling down")

        driver.execute_script("window.scrollBy(0, 200)")
        driver.save_screenshot('screenshot.png')
            
        print("ss will taken")

        im = Image.open("screenshot.png")

        im_cropped = im.crop((0, 15, 1090, 600))

        im_cropped.save("newimage.png")

        image = Image.open("newimage.png")

        width, height = image.size

        new_width = 924
        new_height = 500

        resized_image = image.resize((new_width, new_height))

        resized_image.save("newimage2.png")

        print("Screenshot cropped")

        driver.quit()

        print("Screenshot successfully downloaded")

    except Exception as e:
        print(f"An error occurred while downloading the screenshot, reason: {e}")

    if review:
        review = extract_numbers(review)


    filtered_info["Price"] = price

    for key in search_keys:
        if key in product_info:
            stripped_product_info[key] = product_info[key].strip()
                
            if ';' in stripped_product_info[key]:
                dimensions, weight_info = stripped_product_info[key].split(';', 1)
                new_key = f"Product Weight"
                    
                filtered_info[new_key] = weight_info.strip()
                    
                filtered_info[key] = dimensions.strip()
            else:
                filtered_info[key] = stripped_product_info[key]
                

    for key, value in filtered_info.items():
        new_key = translation_dict.get(key.strip(), key.strip())
            
        new_value = ' '.join(value_translation_dict.get(word, word) for word in value.split(' '))
            
        new_output[new_key] = new_value


    product_title = product_title.strip()
    product_title = product_title[:100]
    try:
        prs = Presentation()

        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        slide_layout = prs.slide_layouts[1]  
        slide = prs.slides.add_slide(slide_layout)

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if shape.text_frame:
                p = shape.text_frame.text = " "

        title = slide.shapes.title

        title.text = "Example products"

        img_path = "newimage2.png"

        with Image.open(img_path) as img:
            width, height = img.size

        width_inch = width / 96.0
        height_inch = height / 96.0

        left = Inches(0.4)
        top = Inches(1.25)
        slide.shapes.add_picture(img_path, left, top, Inches(width_inch-2), Inches(height_inch))

        txBox = slide.shapes.add_textbox(Inches(9), Inches(1), Inches(3.40), Inches(6.30))

        line = txBox.line
        line.color.rgb = RGBColor(91, 155, 213)
        line.width = Inches(0.005)

        content = slide.placeholders[1]
        text_frame = content.text_frame
        p = text_frame.add_paragraph()
        p.text = " "

        tf = txBox.text_frame
        tf.word_wrap = True

        for key, value in new_output.items():
            if "Ounces" in str(value):
                try:
                    weight_value = float(re.findall(r'\d+', value)[0])
                    new_value = weight_value * 28
                    value = f"{new_value} Gram"
                except (ValueError, IndexError):
                    pass 
            elif "lb" in str(value) or "lbs" in str(value):
                try:
                    weight_value = float(re.findall(r'\d+', value)[0])
                    new_value = weight_value * 454 
                    value = f"{new_value} Gram"
                except (ValueError, IndexError):
                        pass

            p = tf.add_paragraph()
            p.text = f"{key}: {value}"
            p.font.bold = False
            p.font.size = Pt(18)
            p.font.name = 'Calibri Light'
            p.alignment = 1

        p = tf.add_paragraph()
        p.text = f"Review numbers: {review}"
        p.font.bold = False
        p.font.size = Pt(18)
        p.font.name = 'Calibri Light'
        p.alignment = 1

        p = tf.add_paragraph()
        p.text = f"Seller: {seller_name}" 
        p.font.bold = False
        p.font.size = Pt(18)
        p.font.name = 'Calibri Light'
        p.alignment = 1

        tf.margin_top = Inches(1.70)

        prs.save(f"{product_title}.pptx")

        os.remove("newimage.png")
        os.remove("screenshot.png")
        os.remove("newimage2.png")
        print("PowerPoint project created")
    except Exception as e:
        print(f"An error occurred while creating PowerPoint Project, hata:{e}")
        import traceback
        traceback.print_exc()


#yapılacaklar listes:
#This instructs the user to install the necessary Python packages using pip by running the given command in the Terminal.
#pip install requests beautifulsoup4 pillow python-pptx selenium fake_useragent


search_keys = [
    "Produktabmessungen", "Produktabmessungen   ",
    "Ürün Boyutları", "Ürün Boyutları   ",
    "Paket Boyutları", "Paket Boyutları   ",
    "Item Dimensions LxWxH", "Item Dimensions LxWxH   ",
    "Product Dimensions", "Product Dimensions   ",
    "Package Dimensions", "Package Dimensions   ",
    "Verpackungsabmessungen", "Verpackungsabmessungen   ",
    "Ürün Boyutları", "Ürün Boyutları   ",
]
search_weights_key = [    
    "Netto-Gewicht", "Netto-Gewicht   " ,
    "Stückzahl", "Stückzahl   ",
    "Artikelgewicht", "Artikelgewicht   ",
    "Weight", "Weight   ",
    "Units", "Units   ",
    "Item Weight", "Item Weight   "
]
translation_dict = {
    "Produktabmessungen": "Product Dimensions",
    "Item Dimensions LxWxH": "Product Dimensions",
    "Product Dimensions": "Product Dimensions",
    "Verpackungsabmessungen": "Product Dimensions",
    "Paket Boyutları": "Product Dimensions",
    "Ürün Boyutları" : "Product Dimensions",
    "Package Dimensions": "Product Dimensions"
}
translation_of_weights = {   
    "Netto-Gewicht": "Item Weight",
    "Stückzahl": "Item Weight",
    "Artikelgewicht": "Item Weight",
    "Weight": "Item Weight",
    "Units": "Item Weight",
    "Item Weight": "Item Weight",
    "Ürün Ağırlığı": "Item Weight"
}
value_translation_dict = {
    'Gramm': 'Gram',
    'Kilogramm': 'Kilogram',
    'kg': 'Kilogram',
    'g': 'Gram',
    'Ounces': 'Ounces'
}


headers = {
    'authority': 'www.amazon.com',
    'accept': 'text/html,application/xhtml+xml,application/xml;q=0.9,image/avif,image/webp,image/apng,*/*;q=0.8,application/signed-exchange;v=b3;q=0.7',
    'cache-control': 'no-cache',
    'dnt': '1',
    'pragma': 'no-cache',
    'sec-ch-ua': '"Chromium";v="116", "Not)A;Brand";v="24", "Google Chrome";v="116"',
    'sec-ch-ua-mobile': '?0',
    'sec-ch-ua-platform': '"Windows"',
    'sec-fetch-dest': 'document',
    'sec-fetch-mode': 'navigate',
    'sec-fetch-site': 'cross-site',
    'sec-fetch-user': '?1',
    'upgrade-insecure-requests': '1',
    'user-agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/116.0.0.0 Safari/537.36',
}


while True:
    print("Menu")
    print("1: Single info retrieval.")
    print("2: List-based info retrieval.")
    print("3: Exit")

    secim = input("Make your choice: ")
    
    if secim == '1':
        print("1: Single info retrieval selected.")

        url = input("Enter URL: ")
        aylık_satış = input(f"{url} \nLinki yukarda ki eşya için aylık satış sayısını girin: ")

        the_machine(url, aylık_satış)

    elif secim == '2':
        print("2: List-based info retrieval selected.")

        Tk().withdraw()
        filepath = askopenfilename(filetypes = (("Text files", "*.txt"), ("all files", "*.*")))

        if filepath:
            with open(filepath, 'r') as f:
                urls = f.readlines()
            for url in urls: 
                url = url.strip()

                the_machine(url)

    elif secim == '3':
        print("Exiting...")
        break 
    else:
        print("Invalid selection.")
